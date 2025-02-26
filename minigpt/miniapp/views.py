from django.shortcuts import render, redirect , get_object_or_404
from django.utils.dateparse import parse_datetime
from django.contrib.auth import authenticate, login, logout
from django.contrib.auth.models import User
from django.contrib import messages
from django.http import JsonResponse
import google.generativeai as genai
from PyPDF2 import PdfReader
from pptx import Presentation
import os
from django.views.decorators.csrf import csrf_exempt
from .models import *


#-----------------------------------------Chat Model----------------------------------------------------------
genai.configure(api_key="AIzaSyBbYWGuVstJ7_dAjwvpAfRLesvKOhrxI30")
model = genai.GenerativeModel('gemini-1.5-pro')
chat_session = model.start_chat()

@csrf_exempt
def get_ai_response(message, file_context=""):
    combined_input = f"File Content:\n{file_context}\n\nUser Query:\n{message}"
    chat_session = model.start_chat()
    response = chat_session.send_message(combined_input)
    return response.text
#----------------------------------------------------------------------------------------------------------------

#-------------------------------------------User-------------------------------------------------------------------
def register_view(request):
    if request.method == 'POST':
        username = request.POST.get('username')
        password = request.POST.get('password')
        confirm_password = request.POST.get('confirm_password')
        email = request.POST.get('email')

        if password != confirm_password:
            messages.error(request, 'Passwords do not match.')
            return redirect('register')

        if User.objects.filter(username=username).exists():
            messages.error(request, 'Username already taken.')
            return redirect('register')

        user = User.objects.create_user(username=username, password=password, email=email)
        user.save()
        messages.success(request, 'Registration successful. Please log in.')
        return redirect('login')

    return render(request, 'register.html')


def login_view(request):
    if request.method == 'POST':
        username = request.POST.get('username')
        password = request.POST.get('password')

        user = authenticate(request, username=username, password=password)
        if user is not None:
            login(request, user)
            if user.is_superuser:
                return redirect('adminDashboard')
            return redirect('index')
        else:
            messages.error(request, 'Invalid username or password.')
            return redirect('login')

    return render(request, 'login.html')


def get_chat_history(request):
    chats = ChatHistory.objects.filter(user=request.user).values('user_message', 'ai_response', 'timestamp')
    return JsonResponse(list(chats), safe=False)

#-------------------------------------------------------------------------------------------------------------------

#-----------------------------------------------------Chat Views----------------------------------------------------
@csrf_exempt
def extract_text_from_pdf(file):
    reader = PdfReader(file)
    return "\n".join(page.extract_text() for page in reader.pages if page.extract_text())

@csrf_exempt
def extract_text_from_txt(file):
    return file.read().decode('utf-8')

@csrf_exempt
def extract_text_from_ppt(file):
    prs = Presentation(file)
    return "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))

@csrf_exempt
def process_uploaded_file(file):
    ext = os.path.splitext(file.name)[1].lower()
    if ext == ".pdf":
        return extract_text_from_pdf(file)
    elif ext == ".txt":
        return extract_text_from_txt(file)
    elif ext == ".pptx":
        return extract_text_from_ppt(file)
    else:
        return "Unsupported file type."

@csrf_exempt
def chat_view(request):
    return render(request, 'miniapp/chat.html')


@csrf_exempt
def chat_response(request):
    if request.method == 'POST':
        message = request.POST.get('message')
        uploaded_file = request.FILES.get('uploadedFile')
        file_context = process_uploaded_file(uploaded_file) if uploaded_file else ""

        # Get AI response
        ai_response = get_ai_response(message, file_context)

        # Save chat history
        ChatHistory.objects.create(
            user=request.user,
            user_message=message,
            ai_response=ai_response
        )

        return JsonResponse({'response': ai_response})
    return JsonResponse({'error': 'Invalid request'}, status=400)

#---------------------------------------------------------------------------------------------------------------------

#-----------------------------------------------------Common----------------------------------------------------------
def index(request):
    return render(request, 'index.html')

def logout_view(request):
    logout(request)
    return redirect('login')
#-------------------------------------------------------------------------------------------------------------------

#----------------------------------------------------------Admin---------------------------------------------------
def admin_dashboard_view(request):
    profiles = Profile.objects.select_related('user').all()
    return render(request, 'admin_dashboard.html', {'profiles': profiles})


def view_all_users(request):
    users = User.objects.all().select_related('profile')
    return render(request, 'view_all_users.html', {'users': users})

def update_user(request, user_id):
    user = get_object_or_404(User, id=user_id)
    profile = get_object_or_404(Profile, user=user)

    if request.method == 'POST':
        username = request.POST.get('username')
        email = request.POST.get('email')
        bio = request.POST.get('bio')
        is_premium = request.POST.get('is_premium') == 'on'
        premium_expiry = request.POST.get('premium_expiry')
        profile_picture = request.FILES.get('profile_picture')

        # Update user fields
        user.username = username
        user.email = email
        user.save()

        # Update profile fields
        profile.bio = bio
        profile.is_premium = is_premium
        profile.premium_expiry = parse_datetime(premium_expiry) if premium_expiry else None
        if profile_picture:
            profile.profile_picture = profile_picture
        profile.save()

        messages.success(request, 'User updated successfully.')
        return redirect('viewAllUsers')

    return render(request, 'update_user.html', {'user_obj': user, 'profile': profile})

def delete_user(request, user_id):
    user = get_object_or_404(User, id=user_id)
    if request.method == 'POST':
        user.delete()
        messages.success(request, 'User deleted successfully.')
        return redirect('viewAllUsers')
    return render(request, 'confirm_delete.html', {'user_obj': user})

#-----------------------------------------------------------------------------------------------------------------